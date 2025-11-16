import os
import magic
from django.conf import settings
import PyPDF2
import docx
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract


class DocumentProcessor:
    """Utility class for processing different document types"""

    @staticmethod
    def extract_text_from_file(file_path, file_type):
        """Extract text content from various file types"""
        try:
            if file_type == 'PDF':
                return DocumentProcessor._extract_from_pdf(file_path)
            elif file_type == 'DOCX':
                return DocumentProcessor._extract_from_docx(file_path)
            elif file_type == 'PPTX':
                return DocumentProcessor._extract_from_pptx(file_path)
            elif file_type == 'XLSX':
                return DocumentProcessor._extract_from_excel(file_path)
            elif file_type == 'TXT':
                return DocumentProcessor._extract_from_txt(file_path)
            elif file_type == 'IMAGE':
                return DocumentProcessor._extract_from_image(file_path)
            else:
                return ""
        except Exception as e:
            raise Exception(f"Error extracting text from {file_type}: {str(e)}")

    @staticmethod
    def _extract_from_pdf(file_path):
        """Extract text from PDF files"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()

    @staticmethod
    def _extract_from_docx(file_path):
        """Extract text from Word documents"""
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()

    @staticmethod
    def _extract_from_pptx(file_path):
        """Extract text from PowerPoint presentations"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()

    @staticmethod
    def _extract_from_excel(file_path):
        """Extract text from Excel files"""
        text = ""
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text += f"Sheet: {sheet_name}\n"
            text += df.to_string() + "\n\n"
        return text.strip()

    @staticmethod
    def _extract_from_txt(file_path):
        """Extract text from plain text files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read().strip()

    @staticmethod
    def _extract_from_image(file_path):
        """Extract text from images using OCR"""
        try:
            # Use pytesseract for OCR
            return pytesseract.image_to_string(Image.open(file_path))
        except Exception:
            return "Image file - text extraction requires OCR setup"

    @staticmethod
    def get_file_type(file_path):
        """Determine file type using python-magic"""
        mime = magic.Magic(mime=True)
        mime_type = mime.from_file(file_path)

        mime_map = {
            'application/pdf': 'PDF',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'DOCX',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'PPTX',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'XLSX',
            'text/plain': 'TXT',
            'text/markdown': 'MD',
            'image/jpeg': 'IMAGE',
            'image/png': 'IMAGE',
            'image/gif': 'IMAGE',
        }

        return mime_map.get(mime_type, 'OTHER')